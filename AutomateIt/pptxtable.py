#!/usr/bin/env python

######################################################################
# @author      : Linus Fernandes (linusfernandes at gmail dot com)
# @file        : pptxtable
# @created     : Monday Jan 09, 2023 21:51:05 IST
# @description :
# -*- coding: utf-8 -*-'
######################################################################
import collections
import collections.abc

from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes
shapes.title.text = 'Students Data'
rows = 4
cols = 3
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(1.2)
table = shapes.add_table(rows, cols, left, top, width, height).table
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(2.0)
table.columns[2].width = Inches(2.0)
table.cell(0, 0).text = 'Sr. No.'
table.cell(0, 1).text = 'Student Name'
table.cell(0, 2).text = 'Student Id'
students = {1: ["John", 115], 2: ["Mary", 119], 3: ["Alice", 101]}
for i in range(len(students)):
    table.cell(i + 1, 0).text = str(i + 1)
    table.cell(i + 1, 1).text = str(students[i + 1][0])
    table.cell(i + 1, 2).text = str(students[i + 1][1])
prs.save('table.pptx')
