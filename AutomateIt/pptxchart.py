#!/usr/bin/env python
"""
PPTXChart.

######################################################################
# @author      : Linus Fernandes (linusfernandes at gmail dot com)
# @file        : pptxchart
# @created     : Tuesday Jan 10, 2023 10:07:53 IST
# @description :
# -*- coding: utf-8 -*-'
######################################################################
"""
import collections  # noqa F401 #pylint: disable=unused-import
import collections.abc  # noqa F401 #pylint: disable=unused-import


from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.chart.data import ChartData
from pptx.chart.data import XyChartData
from pptx.chart.data import BubbleChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Inches
from pptx.util import Pt

# create presentation with 1 slide ------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# define chart data ---------------------
chart_data = CategoryChartData()
chart_data.categories = ["East", "West", "Midwest"]
chart_data.add_series("Series 1", (19.2, 21.4, 16.7))

# add chart to slide --------------------
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

slide = prs.slides.add_slide(prs.slide_layouts[5])
chart_data = ChartData()
chart_data.categories = ["East", "West", "Midwest"]
chart_data.add_series("Q1 Sales", (19.2, 21.4, 16.7))
chart_data.add_series("Q2 Sales", (22.3, 28.6, 15.2))
chart_data.add_series("Q3 Sales", (20.4, 26.3, 14.2))

graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

chart = graphic_frame.chart

category_axis = chart.category_axis
category_axis.has_major_gridlines = True
category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
category_axis.tick_labels.font.italic = True
category_axis.tick_labels.font.size = Pt(24)

value_axis = chart.value_axis
value_axis.maximum_scale = 50.0
value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
value_axis.has_minor_gridlines = True

tick_labels = value_axis.tick_labels
tick_labels.number_format = '0"%"'
tick_labels.font.bold = True
tick_labels.font.size = Pt(14)

plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels

data_labels.font.size = Pt(13)
data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
data_labels.position = XL_LABEL_POSITION.INSIDE_END

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.RIGHT
chart.legend.include_in_layout = False

slide = prs.slides.add_slide(prs.slide_layouts[5])

chart_data = XyChartData()

series_1 = chart_data.add_series("Model 1")
series_1.add_data_point(0.7, 2.7)
series_1.add_data_point(1.8, 3.2)
series_1.add_data_point(2.6, 0.8)

series_2 = chart_data.add_series("Model 2")
series_2.add_data_point(1.3, 3.7)
series_2.add_data_point(2.7, 2.3)
series_2.add_data_point(1.6, 1.8)

chart = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data).chart

slide = prs.slides.add_slide(prs.slide_layouts[5])

chart_data = BubbleChartData()

series_1 = chart_data.add_series("Series 1")
series_1.add_data_point(0.7, 2.7, 10)
series_1.add_data_point(1.8, 3.2, 4)
series_1.add_data_point(2.6, 0.8, 8)

chart = slide.shapes.add_chart(XL_CHART_TYPE.BUBBLE, x, y, cx, cy, chart_data).chart

slide = prs.slides.add_slide(prs.slide_layouts[5])

chart_data = ChartData()
chart_data.categories = ["Q1 Sales", "Q2 Sales", "Q3 Sales"]
chart_data.add_series("West", (32.2, 28.4, 34.7))
chart_data.add_series("East", (24.3, 30.6, 20.2))
chart_data.add_series("Midwest", (20.4, 18.3, 26.2))

x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart

chart.has_legend = True
chart.legend.include_in_layout = False
chart.series[0].smooth = True

slide = prs.slides.add_slide(prs.slide_layouts[5])
chart_data = ChartData()
chart_data.categories = ["West", "East", "North", "South", "Other"]
chart_data.add_series("Series 1", (0.135, 0.324, 0.180, 0.235, 0.126))

chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

chart.plots[0].has_data_labels = True
data_labels = chart.plots[0].data_labels
data_labels.number_format = "0%"
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

prs.save("chart.pptx")
