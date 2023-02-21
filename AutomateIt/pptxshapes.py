#!/usr/bin/env python
"""
PPTXShapes.

######################################################################
# @author      : Linus Fernandes (linusfernandes at gmail dot com)
# @file        : pptxshapes
# @created     : Monday Jan 09, 2023 21:19:48 IST
# @description :
# -*- coding: utf-8 -*-'
######################################################################
"""
import collections  # noqa F401 #pylint: disable=unused-import
import collections.abc  # noqa F401 #pylint: disable=unused-import

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE
from pptx.util import Inches

prs = Presentation()
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes
shapes.title.text = "Adding Shapes"
shape1 = shapes.add_shape(
    MSO_SHAPE.RECTANGULAR_CALLOUT,  # pylint: disable=no-member
    Inches(3.5),
    Inches(2),
    Inches(2),
    Inches(2),
)
shape1.fill.solid()
shape1.fill.fore_color.rgb = RGBColor(0x1E, 0x90, 0xFF)
shape1.fill.fore_color.brightness = 0.4
shape1.text = "See! There is home!"
shape2 = shapes.add_shape(
    MSO_SHAPE.ACTION_BUTTON_HOME,  # pylint: disable=no-member
    Inches(3.5),
    Inches(5),
    Inches(2),
    Inches(2),
)
shape2.text = "Home"
prs.save("shapes.pptx")
