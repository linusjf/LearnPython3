#!/usr/bin/env python
"""
PPTXPuctures.

######################################################################
# @author      : Linus Fernandes (linusfernandes at gmail dot com)
# @file        : pptxpictures
# @created     : Tuesday Jan 10, 2023 09:39:43 IST
# @description :
# -*- coding: utf-8 -*-'
######################################################################
"""
import collections  # noqa F401 #pylint: disable=unused-import
import collections.abc  # noqa F401 #pylint: disable=unused-import
from pptx import Presentation
from pptx.util import Inches

IMG_PATH = "python.png"
IMG_PATH2 = "learn_python.jpeg"
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
left = top = Inches(2)
pic = slide.shapes.add_picture(IMG_PATH, left, top, height=Inches(2), width=Inches(3))
left = Inches(2)
top = Inches(5)
height = Inches(2)
pic = slide.shapes.add_picture(IMG_PATH2, left, top, height=height)
prs.save("picture.pptx")
