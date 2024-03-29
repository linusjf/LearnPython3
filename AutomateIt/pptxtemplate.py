#!/usr/bin/env python
"""
PPTXTemplate.

######################################################################
# @author      : Linus Fernandes (linusfernandes at gmail dot com)
# @file        : pptxtemplate
# @created     : Monday Jan 09, 2023 12:03:19 IST
# @description :
# -*- coding: utf-8 -*-'
######################################################################
"""
import collections  # noqa F401 #pylint: disable=unused-import
import collections.abc  # noqa F401 #pylint: disable=unused-import
from pptx import Presentation  # type: ignore

prs = Presentation("Sample_ppt.pptx")
first_slide = prs.slides[0]
first_slide.shapes[0].text_frame.paragraphs[0].text = "Hello!"
slide = prs.slides.add_slide(prs.slide_layouts[1])
text_frame = slide.shapes[0].text_frame
p = text_frame.paragraphs[0]
p.text = "This is a paragraph"
prs.save("new_ppt.pptx")
