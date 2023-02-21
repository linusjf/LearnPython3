#!/usr/bin/env python
"""
PPXTextBox.

######################################################################
# @author      : Linus Fernandes (linusfernandes at gmail dot com)
# @file        : pptxtextbox
# @created     : Monday Jan 09, 2023 14:05:28 IST
# @description :
# -*- coding: utf-8 -*-'
######################################################################
"""
import collections  # noqa F401 #pylint: disable=unused-import
import collections.abc  # noqa F401 #pylint: disable=unused-import

from pptx import Presentation  # type: ignore
from pptx.util import Inches, Pt  # type: ignore

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
txbox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(5), Inches(1))
tf = txbox.text_frame
tf.text = "Wow! I'm inside a textbox"
p = tf.add_paragraph()
p.text = "Adding a new text"
p.font.bold = True
p.font.italic = True
p.font.size = Pt(30)
prs.save("textBox.pptx")
