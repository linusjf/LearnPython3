#!/usr/bin/env python
"""
PPTXSalesData.

######################################################################
# @author      : Linus Fernandes (linusfernandes at gmail dot com)
# @file        : pptxsalesdata
# @created     : Tuesday Jan 10, 2023 16:17:58 IST
# @description :
# -*- coding: utf-8 -*-'
######################################################################
"""
import collections  # noqa F401 #pylint: disable=unused-import
import collections.abc  # noqa F401 #pylint: disable=unused-import

from datetime import datetime
import pandas as pd  # type: ignore  # pylint: disable=import-error
from pptx import Presentation  # type: ignore
from pptx.chart.data import ChartData  # type: ignore
from pptx.enum.chart import XL_CHART_TYPE  # type: ignore
from pptx.enum.chart import XL_LEGEND_POSITION  # type: ignore
from pptx.util import Inches  # type: ignore

xls_file = pd.ExcelFile("Sales_Data.xlsx")
prs = Presentation("sample_ppt.pptx")
first_slide = prs.slides[0]
first_slide.shapes[0].text_frame.paragraphs[
    0
].text = f"Weekly Sales Report {datetime.now().strftime('%D')}"
first_slide.placeholders[1].text = "Author: Alex, alex@innova8"
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
slide.shapes.title.text = "% Revenue for Accounts"
df = xls_file.parse("Sales")
df["total"] = df["Items"] * df["UnitPrice"]
plot = df.groupby("Account")["total"].sum().plot(kind="pie", autopct="%.2f", fontsize=20)
f = plot.get_figure()
f.savefig("result.png", bbox_inches="tight", dpi=400)
left = Inches(2.5)
top = Inches(3)
pic = slide.shapes.add_picture("result.png", left, top, height=Inches(4), width=Inches(5))
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.shapes.title.text = "Sales Manager Performance"
df = xls_file.parse("Sales")
df["total"] = df["Items"] * df["UnitPrice"]
mgr_data = df.groupby(["Manager"])["total"].sum()
managers = mgr_data.index.values.tolist()
sales = []
for mgr in managers:
    sales.append(mgr_data.loc[mgr])

chart_data = ChartData()
chart_data.categories = managers
chart_data.add_series("Series 1", tuple(sales))
x, y, cx, cy = Inches(2), Inches(3), Inches(6), Inches(4)
chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
prs.save("sales.pptx")
