#!/usr/bin/env python
"""
PpTxTwoContent.

######################################################################
# @author      : Linus Fernandes (linusfernandes at gmail dot com)
# @file        : pptxtwocontent
# @created     : Monday Jan 09, 2023 13:29:52 IST
# @description :
# -*- coding: utf-8 -*-'
######################################################################
"""
import collections  # noqa F401 #pylint: disable=unused-import
import collections.abc  # noqa F401 #pylint: disable=unused-import

from pptx import Presentation  # type: ignore

prs = Presentation()
two_content_slide_layout = prs.slide_layouts[3]
slide = prs.slides.add_slide(two_content_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = "Adding a Two Content Slide"
body_shape = shapes.placeholders[1]
tf = body_shape.text_frame
tf.text = "This is line 1."
p = tf.add_paragraph()
p.text = "Again a Line 2.."
p.level = 1
p = tf.add_paragraph()
p.text = "And this is line 3..."
p.level = 2
prs.save("twocontent.pptx")
