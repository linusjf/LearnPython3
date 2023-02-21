#!/usr/bin/env python
"""
PPTXWrite.

######################################################################
# @author      : Linus Fernandes (linusfernandes at gmail dot com)
# @file        : pptxwrite
# @created     : Sunday Jan 08, 2023 12:31:53 IST
# @description :
# -*- coding: utf-8 -*-'
######################################################################
"""
import collections  # noqa # pylint: disable = unused-import
from collections import abc  # noqa # pylint: disable = unused-import
from pptx import Presentation  # type: ignore

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Yo, Python!"
slide.placeholders[1].text = "Yes it is really awesome"
prs.save("yoPython.pptx")
