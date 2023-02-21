#!/usr/bin/env python
"""
PPTXRead.

######################################################################
# @author      : Linus Fernandes (linusfernandes at gmail dot com)
# @file        : pptxread
# @created     : Sunday Jan 08, 2023 12:12:23 IST
# @description :
# -*- coding: utf-8 -*-'
######################################################################
"""
import collections  # noqa F401 #pylint: disable=unused-import
import collections.abc  # noqa F401 #pylint: disable=unused-import


from pptx import Presentation

PATH_TO_PRESENTATION = "samplepptx.pptx"
prs = Presentation(PATH_TO_PRESENTATION)
print("Presentation object for myprofile file: ", prs)
print("Slides are:")
for slide in prs.slides:
    print("Slide object:", slide)
print("Slide has following objects:")
slide1, slide2 = prs.slides[0], prs.slides[1]
print("Slide Ids: \n", slide1.slide_id, ",", slide2.slide_id)
print("Slide Open XML elements: \n", slide1.element, ",", slide2.element)
print("Slide layouts: \n", slide1.slide_layout.name, ",", slide2.slide_layout.name)  # noqa F501

print("Shapes in the slides")
i = 1
for slide in prs.slides:
    print("Slide", i)
    for shape in slide.shapes:
        print("Shape: ", shape.shape_type)
    i += 1

text_runs = []
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)

print("Text is: ", text_runs)
